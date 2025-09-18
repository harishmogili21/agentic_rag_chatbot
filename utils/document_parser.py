import os
import pandas as pd
import pypdf
import docx
import pptx
from typing import List

def parse_txt(file_path: str) -> str:
    """Parses a text or markdown file."""
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def parse_pdf(file_path: str) -> str:
    """Parses a PDF file."""
    text = ""
    with open(file_path, 'rb') as f:
        reader = pypdf.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text

def parse_docx(file_path: str) -> str:
    """Parses a DOCX file."""
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_pptx(file_path: str) -> str:
    """Parses a PPTX file."""
    pres = pptx.Presentation(file_path)
    text = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def parse_csv(file_path: str) -> str:
    """Parses a CSV file into a string representation."""
    df = pd.read_csv(file_path)
    return df.to_string()

def parse_document(file_path: str) -> str:
    """
    A factory function that parses a document based on its file extension.
    """
    _, extension = os.path.splitext(file_path)
    extension = extension.lower()

    if extension == ".pdf":
        return parse_pdf(file_path)
    elif extension == ".docx":
        return parse_docx(file_path)
    elif extension == ".pptx":
        return parse_pptx(file_path)
    elif extension == ".csv":
        return parse_csv(file_path)
    elif extension in [".txt", ".md"]:
        return parse_txt(file_path)
    else:
        raise ValueError(f"Unsupported file format: {extension}")